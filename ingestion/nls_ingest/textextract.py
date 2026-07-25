"""Per-file text extraction with an OCR fallback for scanned PDFs.

Returns an ExtractResult with the text, page count, and the method used, or
method='needs_ocr' when a PDF yielded too little text and OCR was not enabled.
"""

from __future__ import annotations

import logging
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path

from . import config

# pdfminer (under pdfplumber) logs benign color-operator warnings per page;
# silence everything below error.
logging.getLogger("pdfminer").setLevel(logging.ERROR)
logging.getLogger("pdfplumber").setLevel(logging.ERROR)


@dataclass
class ExtractResult:
    text: str
    pages: int
    method: str  # 'pdf_text' | 'pdf_ocr' | 'docx' | 'textutil' | 'plain' | 'needs_ocr' | 'skipped'


def _extract_pdf_text(path: Path) -> tuple[str, int]:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        n = len(pdf.pages)
        for page in pdf.pages:
            parts.append(page.extract_text() or "")
    return "\n\n".join(parts), n


def _ocr_pdf(path: Path) -> tuple[str, int]:
    """OCR every page via pdf2image + pytesseract. Slow; only when enabled."""
    import pytesseract
    from pdf2image import convert_from_path

    parts: list[str] = []
    with tempfile.TemporaryDirectory() as tmp:
        images = convert_from_path(path, dpi=200, output_folder=tmp, fmt="png")
        for img in images:
            parts.append(pytesseract.image_to_string(img))
    return "\n\n".join(parts), len(parts)


def _textutil(path: Path) -> str:
    """macOS built-in converter for .doc/.rtf/.docx/.html -> plain text."""
    out = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", str(path)],
        capture_output=True, text=True, timeout=300,
    )
    return out.stdout if out.returncode == 0 else ""


def extract(path: Path, ocr: bool = False) -> ExtractResult:
    ext = path.suffix.lower()

    if ext == ".pdf":
        try:
            text, pages = _extract_pdf_text(path)
        except Exception:
            text, pages = "", 0
        if len(text.strip()) >= config.MIN_CHARS_FOR_TEXT:
            return ExtractResult(text, pages, "pdf_text")
        # Too little text -> scanned/image PDF.
        if ocr:
            try:
                text, pages = _ocr_pdf(path)
                return ExtractResult(text, pages, "pdf_ocr")
            except Exception:
                return ExtractResult("", pages, "needs_ocr")
        return ExtractResult(text, pages, "needs_ocr")

    if ext == ".docx":
        try:
            import docx
            doc = docx.Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
            return ExtractResult(text, 0, "docx")
        except Exception:
            text = _textutil(path)
            return ExtractResult(text, 0, "textutil")

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            parts.append(" | ".join(c.text for c in row.cells))
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    parts.append(slide.notes_slide.notes_text_frame.text)
            n = len(prs.slides._sldIdLst)
            return ExtractResult("\n".join(parts), n, "pptx")
        except Exception:
            return ExtractResult("", 0, "skipped")

    if ext in (".doc", ".rtf"):
        return ExtractResult(_textutil(path), 0, "textutil")

    if ext in (".txt", ".md"):
        try:
            return ExtractResult(path.read_text(errors="replace"), 0, "plain")
        except Exception:
            return ExtractResult("", 0, "skipped")

    return ExtractResult("", 0, "skipped")
